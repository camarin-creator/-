from pptx import Presentation
from pptx.util import Inches
import re
import openpyxl
import os
import glob

# 추출 함수 정의 (변경 없음)
def extract_text_from_specific_location(pptx_path, slide_indices, x_min_inch, y_min_inch, x_max_inch, y_max_inch):
    """
    PPT 파일의 특정 슬라이드에서 지정된 위치 범위 내의 텍스트를 추출합니다.
    """
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"오류: '{pptx_path}' 파일을 열 수 없습니다. 자세한 오류: {e}")
        return None 

    all_extracted_texts = {}

    for slide_index in slide_indices:
        if slide_index < 0 or slide_index >= len(prs.slides):
            print(f"경고: '{pptx_path}' 파일의 슬라이드 인덱스 {slide_index + 1}는(은) 유효하지 않습니다. 건너뜜.")
            continue

        slide = prs.slides[slide_index]
        extracted_texts_on_slide = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame and shape.text_frame.text.strip():
                
                shape_x1 = shape.left.inches
                shape_y1 = shape.top.inches
                shape_x2 = shape_x1 + shape.width.inches
                shape_y2 = shape_y1 + shape.height.inches

                if (x_min_inch < shape_x2 and x_max_inch > shape_x1 and
                    y_min_inch < shape_y2 and y_max_inch > shape_y1):
                    
                    text_content = shape.text_frame.text.strip()
                    if text_content:
                        # '#숫자' 패턴을 찾아서 빈 문자열로 대체 (삭제)
                        cleaned_text_content = re.sub(r'#\d+', '', text_content)
                        # '강사: ' 문자열을 찾아서 빈 문자열로 대체 (삭제)
                        cleaned_text_content = re.sub(r'강사:\s*', '', cleaned_text_content)
                        
                        ## 이 부분을 추가하세요 ##
                        # '성우: ' 문자열을 찾아서 빈 문자열로 대체 (삭제)
                        # '성우:' 뒤에 공백이 있을 수 있으니 \s*를 추가
                        cleaned_text_content = re.sub(r'성우:\s*', '', cleaned_text_content)
                        ## 추가 끝 ##

                        #  (수직 탭) 문자를 공백으로 대체하거나 완전히 제거
                        cleaned_text_content = re.sub(r'\u000b|\u000c', ' ', cleaned_text_content).strip()
                        
                        if cleaned_text_content: 
                            extracted_texts_on_slide.append(cleaned_text_content)
                            print(f"  찾음: '{cleaned_text_content}' (원래: '{text_content}', 위치: X:{shape_x1:.2f}~{shape_x2:.2f}, Y:{shape_y1:.2f}~{shape_y2:.2f})")
            
        all_extracted_texts[slide_index + 1] = extracted_texts_on_slide

    return all_extracted_texts

# --- 메인 실행 로직 ---
try:
    current_directory = os.getcwd() 

    ppt_files = glob.glob(os.path.join(current_directory, '*.pptx'), recursive=False) + \
                  glob.glob(os.path.join(current_directory, '*.PPTX'), recursive=False) + \
                  glob.glob(os.path.join(current_directory, '*.ppt'), recursive=False) + \
                  glob.glob(os.path.join(current_directory, '*.PPT'), recursive=False)

    if not ppt_files:
        print(f"\n오류: 현재 폴더('{current_directory}')에 .pptx 또는 .ppt 파일이 없습니다. 파일을 확인해주세요.")
        input("작업이 완료되었습니다. 엔터 키를 누르면 창이 닫힙니다...")
        exit() 

    target_slide_numbers = [1, 2, 3, 4, 5] 

    x_start_location = 0.000  
    y_start_location = 5.673  
    x_end_location = 10.1     
    y_end_location = 7.005    

    print(f"\n--- 현재 폴더('{current_directory}')에서 PPT 파일 찾기 시작 ---")
    print(f"총 {len(ppt_files)}개의 PPT 파일을 찾았습니다.")

    for ppt_file_path in ppt_files:
        current_ppt_filename = os.path.basename(ppt_file_path)
        print(f"\n\n===== '{current_ppt_filename}' 파일 처리 시작 =====")

        base_name = os.path.splitext(current_ppt_filename)[0]
        output_excel_file = f"{base_name}.xlsx"

        wb = openpyxl.Workbook()
        ws = wb.active 
        ws.title = "Extracted Texts" 
        ws.append(["슬라이드 번호", "추출된 텍스트"]) 

        try:
            found_treasures = extract_text_from_specific_location(
                ppt_file_path,
                target_slide_numbers,
                x_start_location,
                y_start_location,
                x_end_location,
                y_end_location
            )

            if found_treasures is None: 
                print(f"오류: '{current_ppt_filename}' 파일을 열 수 없어서 처리할 수 없습니다.")
                ws.append(["오류", "PPT 파일을 열 수 없음 (손상 또는 형식 문제)"])
            elif found_treasures:
                for slide_num, texts in found_treasures.items():
                    print(f"\n[슬라이드 {slide_num}에서 찾은 보물들]:")
                    if texts:
                        for text in texts:
                            print(f"- {text}")
                            ws.append([slide_num, text]) 
                    else:
                        print(f"- 슬라이드 {slide_num}: 지정된 구역에서 텍스트를 찾을 수 없습니다.")
                        ws.append([slide_num, "텍스트 없음"]) 
            else: 
                print(f"'{current_ppt_filename}'에서 지정된 슬라이드에 텍스트를 찾을 수 없습니다.")
                ws.append(["정보", "지정된 슬라이드에 텍스트를 찾을 수 없음"])
            
            wb.save(output_excel_file)
            print(f"\n--- '{output_excel_file}' 엑셀 파일 저장 완료 ---")

        except Exception as file_e: 
            print(f"\n** 심각 오류: '{current_ppt_filename}' 파일 처리 중 예상치 못한 오류 발생 **")
            print(f"오류 내용: {file_e}")
            print(f"이 파일에 대한 엑셀 파일은 생성되지 않거나 오류 정보만 포함될 수 있습니다.")
            ws.append(["처리 오류", f"처리 중 오류 발생: {file_e}"])
            try:
                wb.save(f"{base_name}_ERROR.xlsx") 
                print(f"오류 정보가 포함된 '{base_name}_ERROR.xlsx' 파일이 생성되었습니다.")
            except Exception as save_e:
                print(f"오류 정보를 포함한 엑셀 파일 저장 실패: {save_e}")
                print(f"다른 프로그램에서 '{output_excel_file}' 또는 '{base_name}_ERROR.xlsx' 파일을 열어두었는지 확인해주세요.")


        print(f"\n===== '{current_ppt_filename}' 파일 처리 완료 =====")

    print(f"\n\n--- 모든 PPT 파일 처리 완료 ---")

except Exception as main_e: 
    print(f"\n\n--- 심각한 초기 오류 발생 ---")
    print(f"프로그램 시작 또는 파일 검색 중 예상치 못한 오류가 발생했습니다: {main_e}")
    print(f"이 오류는 초기 설정 또는 파일 접근 문제일 수 있습니다.")

input("모든 작업이 완료되었습니다. 엔터 키를 누르면 창이 닫힙니다...")